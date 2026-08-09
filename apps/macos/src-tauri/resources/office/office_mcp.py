#!/usr/bin/env python3
"""Local Office MCP tools for Bob Work Microsoft plugins.

Runs in a local Python sandbox (python-docx, openpyxl, python-pptx when available).
Falls back to OOXML ZIP inspection when libraries are missing.
"""

from __future__ import annotations

import json
import os
import sys
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

OFFICE_KIND = os.environ.get("BOB_OFFICE_KIND", "documents").strip().lower()
NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}


def send(request_id, result):
    payload = {"jsonrpc": "2.0", "id": request_id, "result": result}
    print(json.dumps(payload, ensure_ascii=False), flush=True)


def send_error(request_id, message: str, code: int = -32000):
    print(
        json.dumps(
            {
                "jsonrpc": "2.0",
                "id": request_id,
                "error": {"code": code, "message": message},
            },
            ensure_ascii=False,
        ),
        flush=True,
    )


def tool_result(data) -> dict:
    text = json.dumps(data, ensure_ascii=False, indent=2)
    return {
        "content": [{"type": "text", "text": text}],
        "structuredContent": data,
        "isError": False,
    }


def require_path(arguments: dict) -> Path:
    raw = arguments.get("path")
    if not isinstance(raw, str) or not raw.strip():
        raise ValueError("path is required")
    path = Path(raw).expanduser()
    if not path.exists():
        raise ValueError(f"File not found: {path}")
    if not path.is_file():
        raise ValueError(f"Not a file: {path}")
    return path


def word_tools() -> list[dict]:
    return [
        {
            "name": "inspect_docx",
            "description": "Inspect a DOCX package: paragraphs, tables, styles metadata.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        },
        {
            "name": "extract_docx_text",
            "description": "Extract plain text from a DOCX file for analysis.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}, "max_paragraphs": {"type": "integer"}},
                "required": ["path"],
            },
        },
        {
            "name": "validate_docx",
            "description": "Validate that a path is a readable DOCX package.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        },
    ]


def excel_tools() -> list[dict]:
    return [
        {
            "name": "inspect_xlsx",
            "description": "Inspect an XLSX workbook: sheet names, dimensions, sample cells.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        },
        {
            "name": "read_xlsx_sheet",
            "description": "Read a rectangular range from an XLSX sheet as JSON rows.",
            "inputSchema": {
                "type": "object",
                "properties": {
                    "path": {"type": "string"},
                    "sheet": {"type": "string"},
                    "max_rows": {"type": "integer"},
                    "max_cols": {"type": "integer"},
                },
                "required": ["path"],
            },
        },
        {
            "name": "validate_xlsx",
            "description": "Validate that a path is a readable XLSX package.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        },
    ]


def ppt_tools() -> list[dict]:
    return [
        {
            "name": "inspect_pptx",
            "description": "Inspect a PPTX deck: slide count, titles, notes presence.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        },
        {
            "name": "list_pptx_slides",
            "description": "List slide titles and text snippets from a PPTX file.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}, "max_slides": {"type": "integer"}},
                "required": ["path"],
            },
        },
        {
            "name": "validate_pptx",
            "description": "Validate that a path is a readable PPTX package.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        },
    ]


def document_tools() -> list[dict]:
    return [
        {
            "name": "inspect_document",
            "description": "Inspect a local document (txt, md, pdf, docx, rtf) metadata.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}},
                "required": ["path"],
            },
        },
        {
            "name": "extract_document_text",
            "description": "Extract text from supported local document formats.",
            "inputSchema": {
                "type": "object",
                "properties": {"path": {"type": "string"}, "max_chars": {"type": "integer"}},
                "required": ["path"],
            },
        },
    ]


def tools_for_kind() -> list[dict]:
    if OFFICE_KIND == "word":
        return word_tools()
    if OFFICE_KIND == "excel":
        return excel_tools()
    if OFFICE_KIND == "ppt":
        return ppt_tools()
    return document_tools()


def inspect_docx(path: Path) -> dict:
    try:
        from docx import Document  # type: ignore

        doc = Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        tables = len(doc.tables)
        return {
            "kind": "docx",
            "path": str(path),
            "engine": "python-docx",
            "paragraph_count": len(paragraphs),
            "table_count": tables,
            "preview_paragraphs": paragraphs[:8],
        }
    except ImportError:
        with zipfile.ZipFile(path) as archive:
            xml = archive.read("word/document.xml")
        root = ET.fromstring(xml)
        texts = [
            node.text.strip()
            for node in root.findall(".//w:t", NS)
            if node.text and node.text.strip()
        ]
        return {
            "kind": "docx",
            "path": str(path),
            "engine": "ooxml-zip",
            "paragraph_count": len(texts),
            "preview_paragraphs": texts[:8],
        }


def extract_docx_text(path: Path, max_paragraphs: int = 40) -> dict:
    try:
        from docx import Document  # type: ignore

        doc = Document(path)
        chunks = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        engine = "python-docx"
    except ImportError:
        with zipfile.ZipFile(path) as archive:
            xml = archive.read("word/document.xml")
        root = ET.fromstring(xml)
        chunks = [
            node.text.strip()
            for node in root.findall(".//w:t", NS)
            if node.text and node.text.strip()
        ]
        engine = "ooxml-zip"
    return {
        "path": str(path),
        "engine": engine,
        "paragraphs": chunks[: max(1, max_paragraphs)],
    }


def validate_docx(path: Path) -> dict:
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
    ok = "[Content_Types].xml" in names and any(name.startswith("word/") for name in names)
    return {"path": str(path), "valid": ok, "format": "docx"}


def inspect_xlsx(path: Path) -> dict:
    try:
        from openpyxl import load_workbook  # type: ignore

        workbook = load_workbook(path, read_only=True, data_only=False)
        sheets = []
        for name in workbook.sheetnames[:8]:
            sheet = workbook[name]
            sheets.append({"name": name, "max_row": sheet.max_row, "max_column": sheet.max_column})
        workbook.close()
        return {"kind": "xlsx", "path": str(path), "engine": "openpyxl", "sheets": sheets}
    except ImportError:
        with zipfile.ZipFile(path) as archive:
            workbook_xml = archive.read("xl/workbook.xml")
        root = ET.fromstring(workbook_xml)
        sheets = [
            sheet.attrib.get("name", "Sheet")
            for sheet in root.findall(".//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}sheet")
        ]
        return {"kind": "xlsx", "path": str(path), "engine": "ooxml-zip", "sheets": [{"name": s} for s in sheets]}


def read_xlsx_sheet(path: Path, sheet: str | None, max_rows: int, max_cols: int) -> dict:
    from openpyxl import load_workbook  # type: ignore

    workbook = load_workbook(path, read_only=True, data_only=True)
    target = sheet or workbook.sheetnames[0]
    if target not in workbook.sheetnames:
        workbook.close()
        raise ValueError(f"Sheet not found: {target}")
    ws = workbook[target]
    rows = []
    for row_index, row in enumerate(ws.iter_rows(max_row=max_rows, max_col=max_cols, values_only=True), start=1):
        rows.append({"row": row_index, "values": list(row)})
    workbook.close()
    return {"path": str(path), "sheet": target, "rows": rows, "engine": "openpyxl"}


def validate_xlsx(path: Path) -> dict:
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
    ok = "xl/workbook.xml" in names
    return {"path": str(path), "valid": ok, "format": "xlsx"}


def inspect_pptx(path: Path) -> dict:
    try:
        from pptx import Presentation  # type: ignore

        deck = Presentation(path)
        slides = []
        for index, slide in enumerate(deck.slides, start=1):
            title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
            slides.append({"index": index, "title": title})
        return {"kind": "pptx", "path": str(path), "engine": "python-pptx", "slide_count": len(slides), "slides": slides[:12]}
    except ImportError:
        with zipfile.ZipFile(path) as archive:
            slide_files = sorted(name for name in archive.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml"))
        return {"kind": "pptx", "path": str(path), "engine": "ooxml-zip", "slide_count": len(slide_files)}


def list_pptx_slides(path: Path, max_slides: int = 20) -> dict:
    try:
        from pptx import Presentation  # type: ignore

        deck = Presentation(path)
        slides = []
        for index, slide in enumerate(deck.slides, start=1):
            if index > max_slides:
                break
            title = slide.shapes.title.text.strip() if slide.shapes.title and slide.shapes.title.text else ""
            body = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip() and shape.text.strip() != title:
                    body.append(shape.text.strip())
            slides.append({"index": index, "title": title, "body": body[:4]})
        return {"path": str(path), "engine": "python-pptx", "slides": slides}
    except ImportError:
        return inspect_pptx(path)


def validate_pptx(path: Path) -> dict:
    with zipfile.ZipFile(path) as archive:
        names = set(archive.namelist())
    ok = "ppt/presentation.xml" in names
    return {"path": str(path), "valid": ok, "format": "pptx"}


def inspect_document(path: Path) -> dict:
    suffix = path.suffix.lower()
    size = path.stat().st_size
    base = {"path": str(path), "extension": suffix, "size_bytes": size}
    if suffix == ".docx":
        return {**base, **inspect_docx(path)}
    if suffix in {".xlsx", ".xlsm"}:
        return {**base, **inspect_xlsx(path)}
    if suffix == ".pptx":
        return {**base, **inspect_pptx(path)}
    return base


def extract_document_text(path: Path, max_chars: int = 12000) -> dict:
    suffix = path.suffix.lower()
    if suffix == ".docx":
        payload = extract_docx_text(path)
        text = "\n".join(payload.get("paragraphs", []))
    elif suffix in {".txt", ".md", ".markdown", ".rtf"}:
        text = path.read_text(encoding="utf-8", errors="replace")
    else:
        text = path.read_text(encoding="utf-8", errors="replace")
    return {"path": str(path), "text": text[: max(1, max_chars)]}


def call_tool(name: str, arguments: dict) -> dict:
    path = require_path(arguments)
    if name == "inspect_docx":
        return inspect_docx(path)
    if name == "extract_docx_text":
        return extract_docx_text(path, int(arguments.get("max_paragraphs", 40)))
    if name == "validate_docx":
        return validate_docx(path)
    if name == "inspect_xlsx":
        return inspect_xlsx(path)
    if name == "read_xlsx_sheet":
        return read_xlsx_sheet(
            path,
            arguments.get("sheet"),
            int(arguments.get("max_rows", 25)),
            int(arguments.get("max_cols", 12)),
        )
    if name == "validate_xlsx":
        return validate_xlsx(path)
    if name == "inspect_pptx":
        return inspect_pptx(path)
    if name == "list_pptx_slides":
        return list_pptx_slides(path, int(arguments.get("max_slides", 20)))
    if name == "validate_pptx":
        return validate_pptx(path)
    if name == "inspect_document":
        return inspect_document(path)
    if name == "extract_document_text":
        return extract_document_text(path, int(arguments.get("max_chars", 12000)))
    raise ValueError(f"Unknown tool: {name}")


def main() -> None:
    for line in sys.stdin:
        if not line.strip():
            continue
        request = json.loads(line)
        request_id = request.get("id")
        method = request.get("method")
        if request_id is None:
            continue
        if method == "initialize":
            send(
                request_id,
                {
                    "protocolVersion": "2025-06-18",
                    "capabilities": {"tools": {}},
                    "serverInfo": {"name": f"bob-work-office-{OFFICE_KIND}", "version": "1.1.0"},
                },
            )
        elif method == "tools/list":
            send(request_id, {"tools": tools_for_kind()})
        elif method == "tools/call":
            tool_name = request.get("params", {}).get("name")
            arguments = request.get("params", {}).get("arguments", {})
            try:
                send(request_id, tool_result(call_tool(tool_name, arguments)))
            except Exception as error:  # noqa: BLE001
                send_error(request_id, str(error))
        else:
            send_error(request_id, "Method not found", -32601)


if __name__ == "__main__":
    main()
